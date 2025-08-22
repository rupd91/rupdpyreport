
def proc_createpptx(sample_report=None, ref_slide_number=None, output_report_name=None, images=None, titles=None, descriptions=None):
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import os, copy, glob

    cwd = os.path.dirname(os.path.abspath(__file__))
    sample_report = sample_report or os.path.join(cwd, 'sample.pptx')
    ref_slide_number = ref_slide_number or 1
    images = images or sorted(glob.glob("*.png") + glob.glob("*.jpg"))
    titles = titles or []
    descriptions = descriptions or []
    for i in range(len(images)):
        if i >= len(titles): titles.append(os.path.splitext(os.path.basename(images[i]))[0])
        if i >= len(descriptions): descriptions.append("--")
    output_report_name = output_report_name or os.path.basename(os.getcwd()) + ".pptx"

    prs = Presentation(sample_report)
    ref = prs.slides[ref_slide_number - 1]
    layout = ref.slide_layout

    for img, title, desc in zip(images, titles, descriptions):
        slide = prs.slides.add_slide(layout)
        img_pos = {}
        for s in ref.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                with open("tempimg.jpg", "wb") as f: f.write(s.image.blob)
                img_pos = {"tempimg.jpg": (s.left, s.top, s.width, s.height)}
            else:
                slide.shapes._spTree.insert_element_before(copy.deepcopy(s.element), "p:extLst")
        for p, (l, t, w, h) in img_pos.items():
            slide.shapes.add_picture(p, l, t, w, h)
            os.remove(p)
        try: slide.shapes._spTree.remove(slide.shapes.title.element)
        except: pass
        for s in slide.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and s.placeholder_format.type == 1: s.text = title
            if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and "description" in s.text.lower():
                s.text = desc
                for p in s.text_frame.paragraphs:
                    for r in p.runs: r.font.size = Pt(12)
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if pics:
            old = pics[0]
            slide.shapes._spTree.remove(old.element)
            slide.shapes.add_picture(img, old.left, old.top, old.width, old.height)

    s = prs.slides._sldIdLst
    prs.part.drop_rel(s[ref_slide_number - 1].rId)
    del s[ref_slide_number - 1]

    prs.save(output_report_name)
    print(f"Powerpoint Report Generated: {output_report_name}")
    return output_report_name

    
def proc_pptx2pdf(pptx_path='output-slides.pptx'):
    import os
    from pptxtopdf import convert
    pptx_path = os.path.abspath(pptx_path)
    convert(pptx_path, os.path.dirname(pptx_path))
    print(f"PPTX Report Converted to PDF")
    return pptx_path.replace(".pptx", ".pdf")

    